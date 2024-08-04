import argparse
import logging
import os
import subprocess
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.DEBUG,
                    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                    handlers=[
                        logging.FileHandler('app.log'),
                        logging.StreamHandler()
                    ])
logger = logging.getLogger(__name__)

def replace_placeholder_text(pptx_path, placeholder, replacement, output_pptx_path):
    """
    Replace placeholder text in a PPTX file and save to a new file.
    
    Args:
        pptx_path (str): Path to the input PPTX file.
        placeholder (str): Text placeholder to be replaced.
        replacement (str): Text to replace the placeholder with.
        output_pptx_path (str): Path to save the modified PPTX file.
    
    Raises:
        Exception: If there's an error during processing.
    """
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, replacement)
        prs.save(output_pptx_path)
        logger.info(f"Replaced placeholder text and saved to {output_pptx_path}.")
    except Exception as e:
        logger.error(f"Failed to replace placeholder text. Error: {e}")
        raise

def pptx_to_pdf(pptx_path, pdf_path):
    """
    Convert a PPTX file to a PDF using LibreOffice.
    
    Args:
        pptx_path (str): Path to the input PPTX file.
        pdf_path (str): Path to save the converted PDF file.
    
    Raises:
        Exception: If there's an error during conversion.
    """
    try:
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path, '--outdir', os.path.dirname(pdf_path)], check=True)
        generated_pdf = os.path.splitext(pptx_path)[0] + '.pdf'
        os.rename(generated_pdf, pdf_path)
        logger.info(f"Converted {pptx_path} to PDF and saved to {pdf_path}.")
    except Exception as e:
        logger.error(f"Failed to convert PPTX to PDF. Error: {e}")
        raise

def main():
    """
    Main function to handle command-line arguments and execute the text replacement and conversion.
    """
    parser = argparse.ArgumentParser(description="Replace text in a PPTX file and convert it to PDF.")
    parser.add_argument('--pptx', default='template.pptx', help='Path to the PPTX template file (default: template.pptx)')
    parser.add_argument('--placeholder', default='placeholder_text', help='Placeholder text to replace (default: placeholder_text)')
    parser.add_argument('--replacement', default='Replacement Text', help='Text to replace the placeholder with (default: Replacement Text)')
    parser.add_argument('--delete-intermediary', action='store_true', help='Delete the intermediary PPTX file after conversion')
    parser.add_argument('--output-dir', default='output', help='Directory for output files (default: output)')
    
    args = parser.parse_args()

    os.makedirs(args.output_dir, exist_ok=True)

    pptx_filename = os.path.basename(args.pptx)
    output_pptx_path = os.path.join(args.output_dir, f"modified_{pptx_filename}")
    output_pdf_path = os.path.join(args.output_dir, f"modified_{os.path.splitext(pptx_filename)[0]}.pdf")

    try:
        replace_placeholder_text(args.pptx, args.placeholder, args.replacement, output_pptx_path)
        pptx_to_pdf(output_pptx_path, output_pdf_path)
        
        if args.delete_intermediary:
            os.remove(output_pptx_path)
            logger.info(f"Deleted intermediary PPTX file: {output_pptx_path}")
    except Exception as e:
        logger.error(f"An error occurred: {e}")

if __name__ == '__main__':
    main()
